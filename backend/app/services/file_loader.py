import os
import tempfile
from typing import List, Dict
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from fastapi import UploadFile
from google import genai
from google.genai import types
import time
import json

def _get_genai_client():
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set")
    return genai.Client(api_key=api_key)


def extract_media_with_gemini(file_path: str, mime_type: str, filename: str) -> List[Dict]:
    """
    Uses Gemini 1.5 Pro to extract a comprehensive text description/transcript 
    from images, audio, or video files. Handles images differently from video/audio
    for optimal RAG ingestion.
    """
    client = _get_genai_client()
    
    # Upload file to Gemini File API
    # We use a context manager if possible, but genai.Client.files.upload returns a file object
    uploaded_file = client.files.upload(file=file_path, config={'mime_type': mime_type})
    
    # Wait for the file to finish processing by Google's backend before we query it
    while uploaded_file.state.name == "PROCESSING":
        time.sleep(2)
        uploaded_file = client.files.get(name=uploaded_file.name)
        
    if uploaded_file.state.name == "FAILED":
        client.files.delete(name=uploaded_file.name)
        raise RuntimeError("Google Gemini failed to process the uploaded media file.")
    
    if mime_type.startswith("image/"):
        prompt = (
            "You are an expert image analyst and explainer for a Retrieval-Augmented Generation (RAG) system. "
            "Analyze this image comprehensively and explain it in extreme detail. Include: "
            "1. The main subject and overall context. "
            "2. A complete transcript of any text present (OCR). "
            "3. Important visual elements, colors, and spatial layout. "
            "4. Detailed data points if it contains charts, graphs, or diagrams. "
            "5. A conceptual explanation of what the image represents or its likely purpose. "
            "You MUST output exactly valid JSON in the following format: "
            '[{"timestamp": "Image", "text": "Detailed explanation and OCR here..."}]'
        )
    else: # Video or Audio
        prompt = (
            "You are an expert video and audio analyst for a Retrieval-Augmented Generation (RAG) system. "
            "Analyze this media comprehensively and break it down into a chronological sequence of scenes or segments. "
            "For each segment, provide: "
            "1. A highly detailed description of all visual actions, objects, scene changes, and on-screen text (if applicable). "
            "2. A complete and accurate transcript of the spoken audio. "
            "You MUST output exactly valid JSON in the following format: "
            '[{"timestamp": "00:00-00:30", "text": "Visuals: [Detailed visual description]\\nTranscript: [Spoken words]"}, ...]'
        )
    
    try:
        # Using gemini-1.5-pro for exceptionally good video tracking and image OCR/reasoning
        response = client.models.generate_content(
            model='gemini-1.5-pro',
            contents=[uploaded_file, prompt],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
            )
        )
        text = response.text.strip()
    except Exception as e:
        raise RuntimeError(f"Failed to generate content with Gemini API: {str(e)}")
    finally:
        # Always clean up the file from Gemini's storage
        client.files.delete(name=uploaded_file.name)
        
    if not text:
        return []
        
    import json
    try:
        scenes = json.loads(text)
        documents = []
        for i, scene in enumerate(scenes):
            scene_text = scene.get("text", "")
            timestamp = scene.get("timestamp", f"Scene {i+1}")
            if scene_text:
                documents.append({
                    "text": f"[{timestamp}] {scene_text}",
                    "page": timestamp, # Using page for timestamp metadata
                    "source": filename
                })
        return documents
    except json.JSONDecodeError:
        # Fallback if json fails
        return [{
            "text": text,
            "page": "1", 
            "source": filename
        }]


async def extract_text_from_file(file: UploadFile) -> List[Dict]:
    """
    Extract text from file WITH metadata.

    Returns a list of:
    {
        "text": str,
        "page": int | None,
        "source": filename
    }
    """

    filename = file.filename
    lower = filename.lower()

    documents: List[Dict] = []
    
    # Read the file content into memory since we might need it as a file path for Gemini
    content = await file.read()

    # -------- PDF --------
    if lower.endswith(".pdf"):
        from io import BytesIO
        reader = PdfReader(BytesIO(content))

        for page_idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()

            if not text:
                continue

            documents.append({
                "text": text,
                "page": page_idx + 1,  # human-readable page number
                "source": filename,
            })

        return documents

    # -------- DOCX --------
    if lower.endswith(".docx") or lower.endswith(".doc"):
        from io import BytesIO
        doc = Document(BytesIO(content))
        full_text = "\n".join(p.text for p in doc.paragraphs).strip()

        if not full_text:
            return []

        documents.append({
            "text": full_text,
            "page": None,
            "source": filename,
        })

        return documents
        
    # -------- PPTX --------
    if lower.endswith(".pptx"):
        from io import BytesIO
        prs = Presentation(BytesIO(content))
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            slide_text = "\n".join(text_runs).strip()
            if slide_text:
                documents.append({
                    "text": slide_text,
                    "page": i + 1, # Treat slides as pages
                    "source": filename
                })
        return documents
        
    # -------- XLSX --------
    if lower.endswith(".xlsx"):
        from io import BytesIO
        wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
        for i, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            sheet_text = []
            sheet_text.append(f"Sheet Name: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Only add rows that have at least some data
                if any(row_values):
                    sheet_text.append(" | ".join(row_values))
            
            full_sheet_text = "\n".join(sheet_text).strip()
            if len(full_sheet_text) > len(f"Sheet Name: {sheet_name}"):
                documents.append({
                    "text": full_sheet_text,
                    "page": i + 1, # Treat sheets as pages
                    "source": filename
                })
        return documents
        
    # -------- MULTIMEDIA (Images, Audio, Video) --------
    mime_type = file.content_type
    if not mime_type or mime_type == "application/octet-stream":
        import mimetypes
        guessed_type, _ = mimetypes.guess_type(filename)
        if guessed_type:
            mime_type = guessed_type

    if mime_type and (mime_type.startswith("image/") or mime_type.startswith("audio/") or mime_type.startswith("video/")):
        # Write to a temporary file because Gemini File API requires a file path
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as temp_file:
            temp_file.write(content)
            temp_path = temp_file.name
            
        try:
            documents = extract_media_with_gemini(temp_path, mime_type, filename)
            return documents
        finally:
            os.remove(temp_path)

    # -------- UNSUPPORTED --------
    raise ValueError(f"Unsupported file type: {filename} with mime {mime_type}")
